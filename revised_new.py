#!/usr/bin/env python3

import os
import json
import time
import re
import logging
import subprocess
import hashlib
import pickle
import sqlite3
import yaml
import threading
import shutil
import queue
import uuid
from datetime import datetime
from concurrent.futures import ThreadPoolExecutor, as_completed
from pathlib import Path
from typing import Optional, Tuple, Any, List, Dict

# Optional third-party libs used by the original script
try:
    import ollama
except Exception:
    ollama = None

try:
    import pypdf
except Exception:
    pypdf = None

try:
    import docx
except Exception:
    docx = None

try:
    import pytesseract
    from PIL import Image
except Exception:
    pytesseract = None
    Image = None

try:
    from pptx import Presentation
except Exception:
    Presentation = None

# --- CONFIGURATION (defaults) ---
INBOX_DIR = str(Path.home() / "Dropbox" / "_courses")
ORGANIZED_DIR = str(Path.home() / "OrganizedFiles")
OLLAMA_MODEL = "gemma3:270m"
IGNORE_DIRS = ['_Unprocessed', '_Ignored']
MAX_CONTENT_LENGTH = 5000
MAX_WORKERS = 1

CONFIG_PATH = "config.yaml"

# Logging setup - minimal to ensure our logger variable exists
logging.basicConfig(level=logging.INFO, format="%(asctime)s %(levelname)s: %(message)s", force=1)
logger = logging.getLogger("FileAgent")
API_LOG_QUEUE: queue.Queue[Tuple[str, str]] = queue.Queue()
API_LOG_DISPATCHER_STARTED = False


def load_config(config_path: str = CONFIG_PATH) -> None:
    """Load configuration from YAML if present (overrides defaults)."""
    global INBOX_DIR, ORGANIZED_DIR, OLLAMA_MODEL, IGNORE_DIRS, MAX_CONTENT_LENGTH, MAX_WORKERS

    if not os.path.exists(config_path):
        logger.debug("No config file at %s, using defaults.", config_path)
        return

    try:
        with open(config_path, "r", encoding="utf-8") as f:
            data = yaml.safe_load(f) or {}
    except Exception as e:
        logger.warning("Failed reading config %s: %s", config_path, e)
        return

    INBOX_DIR = str(Path(data.get("inbox_dir", INBOX_DIR)))
    ORGANIZED_DIR = str(Path(data.get("organized_dir", ORGANIZED_DIR)))
    OLLAMA_MODEL = data.get("ollama_model", OLLAMA_MODEL)
    IGNORE_DIRS = data.get("ignore_dirs", IGNORE_DIRS)
    MAX_CONTENT_LENGTH = int(data.get("max_content_length", MAX_CONTENT_LENGTH))
    MAX_WORKERS = int(data.get("max_workers", MAX_WORKERS))

    logger.info("Configuration loaded from %s", config_path)


def setup_logging():
    """Configure file and console handlers for agent and API logs."""
    # Agent logger
    agent_logger = logging.getLogger("FileAgent")
    agent_logger.setLevel(logging.INFO)

    agent_logger.propagate = False

    # Avoid adding handlers multiple times in case this is called twice
    if not agent_logger.handlers:
        fmt = logging.Formatter("%(asctime)s %(levelname)-8s: %(message)s", "%Y-%m-%d %H:%M:%S")
        fh = logging.FileHandler("agent.log", encoding="utf-8")
        fh.setFormatter(fmt)
        sh = logging.StreamHandler()
        sh.setFormatter(fmt)
        agent_logger.addHandler(fh)
        agent_logger.addHandler(sh)

    # API logger
    api_logger = logging.getLogger("OllamaAPI")
    api_logger.setLevel(logging.INFO)
    api_logger.propagate = False

    if not api_logger.handlers:
        api_fh = logging.FileHandler("api.log", encoding="utf-8")
        api_fh.setFormatter(logging.Formatter("%(asctime)s --- %(message)s", "%Y-%m-%d %H:%M:%S"))
        api_logger.addHandler(api_fh)

    start_api_log_dispatcher(api_logger)

    return agent_logger, api_logger


def start_api_log_dispatcher(api_logger: logging.Logger) -> None:
    """Ensure a background thread flushes request/response pairs sequentially."""
    global API_LOG_DISPATCHER_STARTED
    if API_LOG_DISPATCHER_STARTED:
        return

    def dispatcher():
        while True:
            try:
                request_log, response_log = API_LOG_QUEUE.get()
            except Exception:
                continue
            api_logger.info(request_log)
            api_logger.info(response_log)
            API_LOG_QUEUE.task_done()

    thread = threading.Thread(target=dispatcher, daemon=True, name="ApiLogDispatcher")
    thread.start()
    API_LOG_DISPATCHER_STARTED = True


def enqueue_api_log(request_log: str, response_log: str) -> None:
    API_LOG_QUEUE.put((request_log, response_log))


class FileClassificationCache:
    """Thread-safe cache for classification results keyed by file MD5."""

    def __init__(self, cache_file: str = "file_cache.pkl"):
        self.cache_file = cache_file
        self.lock = threading.RLock()
        self.logger = logging.getLogger("Cache")
        self.cache: Dict[str, Dict[str, Any]] = self._load_cache()

    def _load_cache(self) -> Dict[str, Dict[str, Any]]:
        if not os.path.exists(self.cache_file):
            return {}
        try:
            with open(self.cache_file, "rb") as f:
                data = pickle.load(f) or {}
            self.logger.info("Loaded cache with %d entries", len(data))
            return data
        except Exception as e:
            self.logger.warning("Failed to load cache (%s). Starting fresh.", e)
            return {}

    def _atomic_write(self, data: Dict[str, Any]) -> None:
        tmp = f"{self.cache_file}.tmp"
        with open(tmp, "wb") as f:
            pickle.dump(data, f)
        os.replace(tmp, self.cache_file)

    def save_cache(self) -> None:
        with self.lock:
            try:
                self._atomic_write(self.cache)
            except Exception as e:
                self.logger.error("Failed to save cache: %s", e)

    def get_file_hash(self, file_path: str) -> Optional[str]:
        try:
            hasher = hashlib.md5()
            with open(file_path, "rb") as f:
                for chunk in iter(lambda: f.read(8192), b""):
                    hasher.update(chunk)
            return hasher.hexdigest()
        except Exception as e:
            self.logger.error("Failed to hash %s: %s", file_path, e)
            return None

    def get(self, file_path: str, filename: str) -> Optional[Dict[str, Any]]:
        h = self.get_file_hash(file_path)
        if not h:
            return None
        with self.lock:
            entry = self.cache.get(h)
            if entry:
                self.logger.info("Cache HIT for %s -> %s", filename, entry.get("category"))
                return entry
            return None

    def set(self, file_path: str, classification_data: Dict[str, Any]) -> None:
        h = self.get_file_hash(file_path)
        if not h or not classification_data:
            return
        with self.lock:
            self.cache[h] = {
                "category": classification_data.get("category"),
                "filename": classification_data.get("filename"),
                "cached_at": datetime.utcnow().isoformat()
            }
            # write immediately; safe for small caches. Could be batched in other designs.
            self.save_cache()
            self.logger.info("Cached classification for %s", classification_data.get("filename"))


class FileDatabase:
    """Thread-safe SQLite wrapper to log file movements and errors."""

    def __init__(self, db_path: str = "file_organization.db"):
        self.db_path = db_path
        self.lock = threading.RLock()
        self.logger = logging.getLogger("Database")
        self.conn = sqlite3.connect(db_path, check_same_thread=False)
        self.conn.row_factory = sqlite3.Row
        self.create_tables()

    def create_tables(self) -> None:
        with self.lock:
            cur = self.conn.cursor()
            cur.execute('''
            CREATE TABLE IF NOT EXISTS file_movements (
                id INTEGER PRIMARY KEY AUTOINCREMENT,
                original_path TEXT NOT NULL,
                original_filename TEXT NOT NULL,
                new_path TEXT,
                new_filename TEXT,
                category TEXT NOT NULL,
                processed_date TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
                file_hash TEXT,
                file_size INTEGER,
                extraction_method TEXT,
                cached BOOLEAN DEFAULT 0,
                processing_time_seconds REAL,
                error_message TEXT
            );
            ''')
            cur.execute('''
            CREATE TABLE IF NOT EXISTS categories (
                id INTEGER PRIMARY KEY AUTOINCREMENT,
                category_name TEXT UNIQUE NOT NULL,
                created_date TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
                file_count INTEGER DEFAULT 0
            );
            ''')
            cur.execute('CREATE INDEX IF NOT EXISTS idx_file_hash ON file_movements(file_hash);')
            cur.execute('CREATE INDEX IF NOT EXISTS idx_category ON file_movements(category);')
            self.conn.commit()
            self.logger.info("Database initialized at %s", self.db_path)

    def log_file_movement(self, original_path: str, new_path: str, category: str,
                          file_hash: Optional[str] = None, cached: bool = False,
                          processing_time: Optional[float] = None,
                          file_size: Optional[int] = None) -> None:
        if file_size is None:
            try:
                probe_path = new_path if new_path and os.path.exists(new_path) else original_path
                file_size = os.path.getsize(probe_path) if probe_path and os.path.exists(probe_path) else None
            except Exception:
                file_size = None

        with self.lock:
            try:
                self.conn.execute('''
                INSERT INTO file_movements
                (original_path, original_filename, new_path, new_filename, category, file_hash, file_size, cached, processing_time_seconds)
                VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?)
                ''', (
                    original_path,
                    os.path.basename(original_path),
                    new_path,
                    os.path.basename(new_path) if new_path else None,
                    category,
                    file_hash,
                    file_size,
                    int(bool(cached)),
                    float(processing_time) if processing_time is not None else None
                ))
                # update category stats
                self.conn.execute('INSERT OR IGNORE INTO categories (category_name) VALUES (?);', (category,))
                self.conn.execute('''
                UPDATE categories
                SET file_count = (
                    SELECT COUNT(*) FROM file_movements WHERE category = ? AND new_path IS NOT NULL
                )
                WHERE category_name = ?
                ''', (category, category))
                self.conn.commit()
            except Exception as e:
                logging.getLogger("Database").error("Failed to log file movement: %s", e)

    def log_error(self, file_path: str, error_message: str, category: str = "ERROR") -> None:
        with self.lock:
            try:
                self.conn.execute('''
                INSERT INTO file_movements (original_path, original_filename, category, error_message)
                VALUES (?, ?, ?, ?)
                ''', (file_path, os.path.basename(file_path), category, str(error_message)))
                self.conn.commit()
            except Exception as e:
                logging.getLogger("Database").error("Failed to log error: %s", e)

    def get_statistics(self) -> Dict[str, Any]:
        with self.lock:
            cur = self.conn.cursor()
            stats = {}
            cur.execute('SELECT COUNT(*) FROM file_movements WHERE new_path IS NOT NULL;')
            stats['total_processed'] = cur.fetchone()[0] or 0

            cur.execute('''
            SELECT category, COUNT(*) AS count FROM file_movements
            WHERE new_path IS NOT NULL
            GROUP BY category ORDER BY count DESC;
            ''')
            stats['by_category'] = cur.fetchall()

            cur.execute('''
            SELECT SUM(CASE WHEN cached = 1 THEN 1 ELSE 0 END) AS cached, COUNT(*) as total
            FROM file_movements WHERE new_path IS NOT NULL;
            ''')
            row = cur.fetchone()
            cached = row[0] or 0
            total = row[1] or 0
            stats['cache_hit_rate'] = (cached / total * 100.0) if total > 0 else 0.0

            cur.execute('SELECT AVG(processing_time_seconds) FROM file_movements WHERE processing_time_seconds IS NOT NULL;')
            stats['avg_processing_time'] = cur.fetchone()[0]
            return stats

    def get_file_history(self, filename: Optional[str] = None):
        with self.lock:
            cur = self.conn.cursor()
            if filename:
                cur.execute('''
                SELECT * FROM file_movements WHERE original_filename = ? OR new_filename = ?
                ORDER BY processed_date DESC
                ''', (filename, filename))
            else:
                cur.execute('SELECT * FROM file_movements ORDER BY processed_date DESC LIMIT 100;')
            return cur.fetchall()

    def close(self) -> None:
        with self.lock:
            self.conn.close()


def retry(exceptions, tries: int = 3, delay: float = 1.0, backoff: int = 2):
    """Retry decorator with exponential backoff."""
    def decorator(func):
        import functools
        @functools.wraps(func)
        def wrapper(*args, **kwargs):
            _tries, _delay = tries, delay
            last_exception = None
            while _tries > 0:
                try:
                    return func(*args, **kwargs)
                except exceptions as e:
                    last_exception = e
                    logger = kwargs.get("logger", logging.getLogger("FileAgent"))
                    logger.warning("Transient error: %s. Retrying in %.1fs...", e, _delay)
                    time.sleep(_delay)
                    _tries -= 1
                    _delay *= backoff
            # final attempt will raise the last exception
            raise last_exception
        return wrapper
    return decorator


def sanitize_category(name: Optional[str]) -> str:
    if not name or not isinstance(name, str):
        return "_Unprocessed"
    safe = re.sub(r'[<>:"/\\|?*\x00-\x1f]+', '', name).strip()
    safe = re.sub(r'\s+', ' ', safe)
    return safe[:100] or "_Unprocessed"


def unique_dest_path(folder: str, filename: str) -> str:
    base, ext = os.path.splitext(filename)
    candidate = filename
    n = 1
    while os.path.exists(os.path.join(folder, candidate)):
        candidate = f"{base}_{n}{ext}"
        n += 1
    return os.path.join(folder, candidate)


def is_tesseract_available() -> bool:
    if pytesseract:
        try:
            _ = pytesseract.get_tesseract_version()
            return True
        except Exception:
            pass
    return shutil.which("tesseract") is not None


def _pdf_ocr_with_pdf2image(pdf_path: str, logger: logging.Logger) -> Optional[str]:
    """Attempt to rasterize PDF pages and OCR each page (requires pdf2image)."""
    try:
        from pdf2image import convert_from_path
    except Exception as e:
        logger.debug("pdf2image not available for OCR fallback: %s", e)
        return None

    try:
        pages = convert_from_path(pdf_path, dpi=200)
        texts = []
        for page_img in pages:
            try:
                txt = pytesseract.image_to_string(page_img, lang='eng') if pytesseract else None
                if txt:
                    texts.append(txt)
            except Exception as e:
                logger.debug("OCR on a page failed: %s", e)
        return "\n".join(texts).strip() or None
    except Exception as e:
        logger.warning("pdf2image conversion failed for %s: %s", pdf_path, e)
        return None


def extract_text(file_path: str, logger: logging.Logger) -> Optional[str]:
    ext = os.path.splitext(file_path)[1].lower()
    text_parts: List[str] = []
    try:
        if ext == '.pdf' and pypdf:
            reader = pypdf.PdfReader(file_path)
            for page in reader.pages:
                page_text = page.extract_text() or ""
                text_parts.append(page_text)
            joined = "\n".join(text_parts).strip()
            if joined:
                return joined
            # fallback to OCR via pdf2image if available + tesseract present
            if is_tesseract_available():
                logger.info("No extracted text in PDF, trying OCR fallback for %s", file_path)
                ocr_text = _pdf_ocr_with_pdf2image(file_path, logger)
                return ocr_text
            return None

        if ext == '.docx' and docx:
            d = docx.Document(file_path)
            return "\n".join(p.text for p in d.paragraphs).strip() or None

        if ext == '.doc':
            # convert if soffice available
            out_dir = os.path.dirname(file_path)
            if not shutil.which("soffice"):
                logger.warning("soffice not found: cannot convert .doc %s", file_path)
                return None
            try:
                subprocess.check_call(['soffice', '--headless', '--convert-to', 'docx', '--outdir', out_dir, file_path],
                                      stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
                base_name = os.path.splitext(os.path.basename(file_path))[0]
                docx_path = os.path.join(out_dir, base_name + '.docx')
                if os.path.exists(docx_path) and docx:
                    d = docx.Document(docx_path)
                    text = "\n".join(p.text for p in d.paragraphs).strip() or None
                    try:
                        os.remove(docx_path)
                    except Exception:
                        pass
                    return text
            except Exception as e:
                logger.warning("soffice conversion failed for %s: %s", file_path, e)
            return None

        if ext in ('.pptx',) and Presentation:
            pres = Presentation(file_path)
            for slide in pres.slides:
                for shape in slide.shapes:
                    if getattr(shape, "has_text_frame", False):
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                text_parts.append(run.text)
            return "\n".join(text_parts).strip() or None

        if ext == '.ppt':
            out_dir = os.path.dirname(file_path)
            if not shutil.which("soffice"):
                logger.warning("soffice not found: cannot convert .ppt %s", file_path)
                return None
            try:
                subprocess.check_call(['soffice', '--headless', '--convert-to', 'pptx', '--outdir', out_dir, file_path],
                                      stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
                base_name = os.path.splitext(os.path.basename(file_path))[0]
                pptx_path = os.path.join(out_dir, base_name + '.pptx')
                if os.path.exists(pptx_path) and Presentation:
                    pres = Presentation(pptx_path)
                    for slide in pres.slides:
                        for shape in slide.shapes:
                            if getattr(shape, "has_text_frame", False):
                                for paragraph in shape.text_frame.paragraphs:
                                    for run in paragraph.runs:
                                        text_parts.append(run.text)
                    try:
                        os.remove(pptx_path)
                    except Exception:
                        pass
                    return "\n".join(text_parts).strip() or None
            except Exception as e:
                logger.warning("soffice conversion failed for %s: %s", file_path, e)
            return None

        if ext in ['.jpg', '.jpeg', '.png']:
            if not is_tesseract_available():
                logger.warning("Tesseract not available; cannot OCR image: %s", file_path)
                return None
            if Image is None:
                logger.warning("Pillow not available; cannot open image for OCR: %s", file_path)
                return None
            try:
                img = Image.open(file_path)
                return pytesseract.image_to_string(img, lang='eng').strip() or None
            except Exception as e:
                logger.warning("Image OCR failed for %s: %s", file_path, e)
                return None

        if ext == '.txt':
            try:
                with open(file_path, 'r', encoding='utf-8') as f:
                    return f.read().strip() or None
            except UnicodeDecodeError:
                with open(file_path, 'r', encoding='latin-1') as f:
                    return f.read().strip() or None

        logger.debug("Unsupported file type for extraction: %s", file_path)
        return None

    except Exception as e:
        logger.warning("Error extracting text from %s: %s", os.path.basename(file_path), e)
        return None


def get_existing_categories(dir_path: str) -> List[str]:
    cats: List[str] = []
    try:
        for item in os.listdir(dir_path):
            item_path = os.path.join(dir_path, item)
            if os.path.isdir(item_path) and item not in IGNORE_DIRS:
                cats.append(item)
    except FileNotFoundError:
        pass
    except Exception as e:
        logger.warning("Error scanning categories in %s: %s", dir_path, e)
    return cats


@retry(Exception, tries=3, delay=1, backoff=2)
def call_ollama_chat(model: str, messages: list, fmt: str = 'json', **kwargs) -> Any:
    """Call ollama.chat with retries (if ollama available)."""
    if ollama is None:
        raise RuntimeError("Ollama client not available (import failed)")
    return ollama.chat(model=model, messages=messages, format=fmt)


def get_classification_from_ollama(
    file_content: str,
    filename: str,
    file_path: str,
    existing_categories: List[str],
    logger: logging.Logger,
) -> Optional[dict]:
    if len(file_content) > MAX_CONTENT_LENGTH:
        # Keep a balanced head/tail preview but never exceed MAX_CONTENT_LENGTH
        head_len = (MAX_CONTENT_LENGTH * 2) // 3
        tail_len = MAX_CONTENT_LENGTH - head_len
        content_preview = (
            file_content[:head_len]
            + "\n... [CONTENT SKIPPED] ...\n"
            + file_content[-tail_len:]
        )
    else:
        content_preview = file_content

    content_preview = ' '.join(content_preview.split())

    prompt = f"""
You are an expert file organization AI agent who prioritizes thematic understanding.

Your task:
- Analyze the filename, current path, and content excerpt to determine the best concept/theme category.

Steps to follow:
1. Let the file's content drive your decision. Only reuse an existing category when the content clearly aligns with that folder's scope; otherwise create a new descriptive CamelCase category.
2. When you do reuse an existing category, copy its name exactly.
3. Treat the current file path as supporting context only. Use it for hints (course name, project, etc.) but never let it override the content-driven choice.
4. Ensure the category name is filesystem-safe (no special characters).

List of existing categories: {str(existing_categories)}
Filename: "{filename}"
Current path: "{file_path}"

##ContentSample##
"{content_preview}"
##EndContentSample##

Respond ONLY with a JSON object like:
{{ 
    "category": "CategoryName" 
}}
"""

    excerpt_for_log = content_preview[:500]
    request_id = uuid.uuid4().hex
    truncated_excerpt = (excerpt_for_log + "...") if len(content_preview) > 500 else excerpt_for_log
    request_log = (
        "REQUEST id=%s filename=%s path=%s excerpt=%s"
        % (
            request_id,
            filename,
            file_path,
            truncated_excerpt,
        )
    )

    try:
        response = call_ollama_chat(
            model=OLLAMA_MODEL,
            messages=[{'role': 'user', 'content': prompt}],
            fmt='json',
            logger=logger
        )
    except Exception as e:
        logger.error("Ollama chat failed: %s", e)
        response_log = f"RESPONSE RAW id={request_id} <call failed: {e}>"
        enqueue_api_log(request_log, response_log)
        return None

    # --- normalize the response into 'content' variable ---
    content = None
    try:
        # Common structures:
        # - response is dict like {'content': '...'} or {'message': {'content': '...'}}
        if isinstance(response, dict):
            if 'message' in response and isinstance(response['message'], dict):
                content = response['message'].get('content')
            else:
                content = response.get('content') or response
        else:
            # response may be an object with attributes .message or .content
            if hasattr(response, 'message'):
                msg = getattr(response, 'message')
                if isinstance(msg, dict):
                    content = msg.get('content')
                elif hasattr(msg, 'content'):
                    content = getattr(msg, 'content')
            if content is None and hasattr(response, 'content'):
                content = getattr(response, 'content')

        # If content is still None, fall back to stringifying the response (safe)
        if content is None:
            content = str(response)

        # If content is a string, try a few parsing strategies
        if isinstance(content, str):
            # 1) Direct JSON
            try:
                parsed = json.loads(content)
                content = parsed
            except Exception:
                # 2) Try to find the first JSON object substring {...}
                start = content.find('{')
                end = content.rfind('}')
                if start != -1 and end != -1 and end > start:
                    candidate = content[start:end+1]
                    try:
                        parsed = json.loads(candidate)
                        content = parsed
                    except Exception:
                        logger.debug("Failed to parse JSON substring: %s", candidate)
                        # leave content as original string for later handling
                else:
                    logger.debug("No JSON-like substring found in Ollama response.")
        response_text = str(content)[:2000]
    except Exception as e:
        logger.warning("Failed to normalize Ollama response (%s): %s", type(response), e)
        response_text = str(response)[:2000]
        response_log = f"RESPONSE RAW id={request_id} <normalization failed: {e}> {response_text}"
        enqueue_api_log(request_log, response_log)
        return None
    else:
        response_log = f"RESPONSE RAW id={request_id} {response_text}"
        enqueue_api_log(request_log, response_log)

    if not isinstance(content, dict):
        logger.warning("Ollama content is not a JSON object/dict.")
        return None

    raw_cat = content.get("category")
    if not raw_cat or not isinstance(raw_cat, str):
        logger.warning("Invalid category returned by Ollama: %s", raw_cat)
        return None

    cat = sanitize_category(raw_cat)
    return {"category": cat}


def move_file(original_path: str, category: str, filename: str, logger: logging.Logger) -> Optional[str]:
    if not category or not filename:
        category = "_Unprocessed"
    safe_category = sanitize_category(category)
    dest_folder = os.path.join(ORGANIZED_DIR, safe_category)
    os.makedirs(dest_folder, exist_ok=True)
    dest_path = unique_dest_path(dest_folder, filename)

    try:
        shutil.move(original_path, dest_path)
        logger.info("Moved: %s -> %s", original_path, dest_path)
        return dest_path
    except Exception as e:
        logger.error("Error moving %s: %s", original_path, e)
        return None


def check_prereqs(logger: logging.Logger) -> bool:
    if ollama is None:
        logger.critical("Ollama client not importable. Ensure ollama is installed and available.")
        return False

    try:
        # If ollama has a list() or ping function, call; otherwise skip
        if hasattr(ollama, "list"):
            _ = ollama.list()
    except Exception as e:
        logger.critical("Could not contact Ollama: %s", e)
        return False

    if not is_tesseract_available():
        logger.warning("Tesseract not found. OCR fallbacks disabled.")

    if not shutil.which("soffice"):
        logger.warning("soffice (LibreOffice) not found. .doc/.ppt conversion disabled.")

    return True


def process_single_file(
    filename: str,
    existing_categories: List[str],
    cache: FileClassificationCache,
    db: FileDatabase,
    logger: logging.Logger,
    categories_lock: threading.Lock
) -> Dict[str, Any]:
    start_time = time.time()
    file_path = os.path.join(INBOX_DIR, filename)
    result = {
        "filename": filename,
        "success": False,
        "category": None,
        "cached": False,
        "error": None,
        "processing_time": 0.0
    }

    logger.info("Processing: %s", filename)

    try:
        cached = cache.get(file_path, filename)
        if cached:
            classification = cached
            result["cached"] = True
            logger.info("Using cached classification for %s", filename)
        else:
            content = extract_text(file_path, logger)
            if not content or content.isspace():
                msg = "Empty content or unsupported file type"
                logger.warning(msg + ": %s", filename)
                result['error'] = msg
                db.log_error(file_path, msg)
                result['processing_time'] = time.time() - start_time
                return result

            classification = get_classification_from_ollama(content, filename, file_path, existing_categories, logger)
            if classification:
                cache.set(file_path, classification)

        if not classification:
            msg = f"AI classification failed for {filename}"
            logger.error(msg)
            result['error'] = msg
            db.log_error(file_path, msg)
            result['processing_time'] = time.time() - start_time
            return result

        category = classification.get("category", "Misc")
        filename_new = classification.get("filename", filename)
        file_hash = cache.get_file_hash(file_path)

        try:
            file_size = os.path.getsize(file_path)
        except Exception:
            file_size = None

        dest_path = move_file(file_path, category, filename_new, logger)
        processing_time = time.time() - start_time

        if dest_path:
            db.log_file_movement(original_path=file_path, new_path=dest_path,
                                 category=category, file_hash=file_hash,
                                 cached=result['cached'], processing_time=processing_time,
                                 file_size=file_size)
            result['success'] = True
            result['category'] = category
            result['processing_time'] = processing_time
            if category not in IGNORE_DIRS:
                with categories_lock:
                    if category not in existing_categories:
                        existing_categories.append(category)
                        logger.info("New category added: %s", category)
            logger.info("Success: %s -> %s (%.2fs)", filename, category, processing_time)
        else:
            err = f"Move failed for {filename}"
            result['error'] = err
            db.log_error(file_path, err, category)
            logger.error(err)
    except Exception as e:
        err = f"Unexpected error processing {filename}: {e}"
        result['error'] = err
        db.log_error(file_path, err)
        logger.exception(err)

    result['processing_time'] = time.time() - start_time
    return result


def main(logger: logging.Logger) -> None:
    logger.info("=" * 60)
    logger.info("FILE ORGANIZATION AI AGENT STARTED")
    logger.info("=" * 60)
    logger.info("Inbox: %s", INBOX_DIR)
    logger.info("Output: %s", ORGANIZED_DIR)
    logger.info("AI Model: %s", OLLAMA_MODEL)
    logger.info("Max Workers: %d", MAX_WORKERS)

    cache = FileClassificationCache()
    db = FileDatabase()

    os.makedirs(INBOX_DIR, exist_ok=True)
    os.makedirs(ORGANIZED_DIR, exist_ok=True)

    existing_categories = get_existing_categories(ORGANIZED_DIR)
    categories_lock = threading.Lock()
    logger.info("Existing categories: %s", existing_categories)

    files_to_process = [
        f for f in os.listdir(INBOX_DIR)
        if os.path.isfile(os.path.join(INBOX_DIR, f)) and not f.startswith('.')
    ]

    if not files_to_process:
        logger.info("No files found to organize.")
        stats = db.get_statistics()
        logger.info("Historical stats: %s", stats)
        db.close()
        return

    logger.info("Found %d files to process", len(files_to_process))

    successful = 0
    failed = 0
    start = time.time()

    with ThreadPoolExecutor(max_workers=MAX_WORKERS) as executor:
        future_to_file = {
            executor.submit(process_single_file, f, existing_categories, cache, db, logger, categories_lock): f
            for f in files_to_process
        }
        for fut in as_completed(future_to_file):
            f_name = future_to_file[fut]
            try:
                res = fut.result()
                if res.get("success"):
                    successful += 1
                else:
                    failed += 1
            except Exception as e:
                logger.error("Thread failed for %s: %s", f_name, e)
                failed += 1

    total_time = time.time() - start
    logger.info("Total files: %d, Successful: %d, Failed: %d", len(files_to_process), successful, failed)
    logger.info("Total time: %.2fs, Avg per file: %.2fs", total_time, total_time / max(1, len(files_to_process)))

    stats = db.get_statistics()
    logger.info("Post-run stats: %s", stats)

    db.close()


if __name__ == "__main__":
    logger, _ = setup_logging()
    load_config()
    if not check_prereqs(logger):
        logger.critical("Pre-run checks failed. Fix issues then retry.")
        raise SystemExit(1)

    try:
    main(logger)
    except KeyboardInterrupt:
        logger.info("Interrupted by user. Exiting.")
    except Exception:
        logger.exception("Unhandled exception in main loop.")
