from __future__ import annotations

import csv
import mimetypes
import os
import re
import shutil
import subprocess
from email import policy
from email.parser import BytesParser
from pathlib import Path
from typing import Callable, Dict, Optional

from .models import ExtractionResult

try:
    import pypdf
except Exception:
    pypdf = None

try:
    import docx
except Exception:
    docx = None

try:
    import pytesseract
    from PIL import Image
except Exception:
    pytesseract = None
    Image = None

try:
    from pptx import Presentation
except Exception:
    Presentation = None

try:
    from openpyxl import load_workbook
except Exception:
    load_workbook = None

try:
    import extract_msg
except Exception:
    extract_msg = None


def sniff_mime_type(file_path: Path) -> Optional[str]:
    mime_type, _ = mimetypes.guess_type(str(file_path))
    return mime_type


def is_tesseract_available() -> bool:
    if pytesseract:
        try:
            _ = pytesseract.get_tesseract_version()
            return True
        except Exception:
            pass
    return shutil.which("tesseract") is not None


def _quality_from_text(text: Optional[str]) -> float:
    if not text:
        return 0.0
    size = len(text.strip())
    if size >= 4000:
        return 0.95
    if size >= 1000:
        return 0.8
    if size >= 250:
        return 0.6
    if size >= 50:
        return 0.35
    return 0.15


def _read_text_file(file_path: Path) -> ExtractionResult:
    for encoding in ("utf-8", "latin-1"):
        try:
            text = file_path.read_text(encoding=encoding).strip()
            return ExtractionResult(
                extracted_text=text or None,
                extraction_quality=_quality_from_text(text),
                extractor_name="text",
                metadata={"encoding": encoding},
            )
        except UnicodeDecodeError:
            continue
    return ExtractionResult(errors=["Could not decode text file"], extractor_name="text")


def _strip_html(html: str) -> str:
    return re.sub(r"<[^>]+>", " ", html or "").strip()


def _extract_pdf(file_path: Path) -> ExtractionResult:
    if pypdf is None:
        return ExtractionResult(errors=["pypdf not available"], extractor_name="pdf")
    try:
        reader = pypdf.PdfReader(str(file_path))
        text_parts = [(page.extract_text() or "") for page in reader.pages]
        joined = "\n".join(text_parts).strip()
        result = ExtractionResult(
            extracted_text=joined or None,
            extraction_quality=_quality_from_text(joined),
            extractor_name="pdf",
            metadata={"pages": len(reader.pages)},
        )
        if joined:
            return result

        if is_tesseract_available():
            try:
                from pdf2image import convert_from_path
            except Exception:
                result.errors.append("pdf2image not available for OCR fallback")
                return result

            ocr_texts = []
            for page_image in convert_from_path(str(file_path), dpi=200):
                txt = pytesseract.image_to_string(page_image, lang="eng") if pytesseract else ""
                if txt:
                    ocr_texts.append(txt)
            ocr_text = "\n".join(ocr_texts).strip()
            result.extracted_text = ocr_text or None
            result.extraction_quality = _quality_from_text(ocr_text)
            result.ocr_used = bool(ocr_text)
        return result
    except Exception as exc:
        return ExtractionResult(errors=[f"PDF extraction failed: {exc}"], extractor_name="pdf")


def _extract_docx(file_path: Path) -> ExtractionResult:
    if docx is None:
        return ExtractionResult(errors=["python-docx not available"], extractor_name="docx")
    try:
        document = docx.Document(str(file_path))
        text = "\n".join(p.text for p in document.paragraphs).strip()
        return ExtractionResult(
            extracted_text=text or None,
            extraction_quality=_quality_from_text(text),
            extractor_name="docx",
        )
    except Exception as exc:
        return ExtractionResult(errors=[f"DOCX extraction failed: {exc}"], extractor_name="docx")


def _convert_office_file(file_path: Path, target_format: str) -> Optional[Path]:
    if not shutil.which("soffice"):
        return None
    out_dir = file_path.parent
    cmd = ["soffice", "--headless", "--convert-to", target_format, "--outdir", str(out_dir), str(file_path)]
    proc = subprocess.run(cmd, stdout=subprocess.PIPE, stderr=subprocess.PIPE, text=True, check=False)
    if proc.returncode != 0:
        return None
    return out_dir / f"{file_path.stem}.{target_format}"


def _extract_doc(file_path: Path) -> ExtractionResult:
    converted = _convert_office_file(file_path, "docx")
    if not converted or not converted.exists():
        return ExtractionResult(errors=["DOC conversion failed"], extractor_name="doc")
    result = _extract_docx(converted)
    result.conversion_used = True
    try:
        converted.unlink()
    except Exception:
        pass
    return result


def _extract_pptx(file_path: Path) -> ExtractionResult:
    if Presentation is None:
        return ExtractionResult(errors=["python-pptx not available"], extractor_name="pptx")
    try:
        presentation = Presentation(str(file_path))
        parts = []
        for slide in presentation.slides:
            for shape in slide.shapes:
                if getattr(shape, "has_text_frame", False):
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            parts.append(run.text)
        text = "\n".join(parts).strip()
        return ExtractionResult(
            extracted_text=text or None,
            extraction_quality=_quality_from_text(text),
            extractor_name="pptx",
            metadata={"slides": len(presentation.slides)},
        )
    except Exception as exc:
        return ExtractionResult(errors=[f"PPTX extraction failed: {exc}"], extractor_name="pptx")


def _extract_ppt(file_path: Path) -> ExtractionResult:
    converted = _convert_office_file(file_path, "pptx")
    if not converted or not converted.exists():
        return ExtractionResult(errors=["PPT conversion failed"], extractor_name="ppt")
    result = _extract_pptx(converted)
    result.conversion_used = True
    try:
        converted.unlink()
    except Exception:
        pass
    return result


def _extract_image(file_path: Path) -> ExtractionResult:
    if not is_tesseract_available():
        return ExtractionResult(errors=["tesseract not available"], extractor_name="image")
    if Image is None or pytesseract is None:
        return ExtractionResult(errors=["image OCR dependencies missing"], extractor_name="image")
    try:
        image = Image.open(str(file_path))
        text = pytesseract.image_to_string(image, lang="eng").strip()
        return ExtractionResult(
            extracted_text=text or None,
            extraction_quality=_quality_from_text(text),
            extractor_name="image",
            ocr_used=True,
        )
    except Exception as exc:
        return ExtractionResult(errors=[f"Image OCR failed: {exc}"], extractor_name="image")


def _extract_xlsx(file_path: Path) -> ExtractionResult:
    if load_workbook is None:
        return ExtractionResult(errors=["openpyxl not available"], extractor_name="xlsx")
    try:
        workbook = load_workbook(filename=str(file_path), read_only=True, data_only=True)
        sheet_names = list(workbook.sheetnames)
        text_parts = []
        for sheet_name in sheet_names:
            worksheet = workbook[sheet_name]
            text_parts.append(f"[Sheet] {sheet_name}")
            for row in worksheet.iter_rows(values_only=True):
                row_values = [str(value).strip() for value in row if value not in (None, "")]
                if row_values:
                    text_parts.append(" | ".join(row_values))
        text = "\n".join(text_parts).strip()
        return ExtractionResult(
            extracted_text=text or None,
            extraction_quality=_quality_from_text(text),
            extractor_name="xlsx",
            metadata={"sheets": sheet_names},
        )
    except Exception as exc:
        return ExtractionResult(errors=[f"XLSX extraction failed: {exc}"], extractor_name="xlsx")


def _extract_xls(file_path: Path) -> ExtractionResult:
    converted = _convert_office_file(file_path, "xlsx")
    if converted and converted.exists():
        result = _extract_xlsx(converted)
        result.conversion_used = True
        try:
            converted.unlink()
        except Exception:
            pass
        return result
    return ExtractionResult(errors=["XLS conversion failed"], extractor_name="xls")


def _extract_eml(file_path: Path) -> ExtractionResult:
    try:
        with open(file_path, "rb") as handle:
            message = BytesParser(policy=policy.default).parse(handle)

        text_parts = []
        metadata = {
            "subject": message.get("subject"),
            "from": message.get("from"),
            "to": message.get("to"),
        }

        if metadata["subject"]:
            text_parts.append(f"Subject: {metadata['subject']}")
        if metadata["from"]:
            text_parts.append(f"From: {metadata['from']}")
        if metadata["to"]:
            text_parts.append(f"To: {metadata['to']}")

        if message.is_multipart():
            for part in message.walk():
                content_type = part.get_content_type()
                disposition = str(part.get("Content-Disposition") or "")
                if "attachment" in disposition.lower():
                    continue
                try:
                    payload = part.get_content()
                except Exception:
                    payload = None
                if not payload:
                    continue
                if content_type == "text/plain":
                    text_parts.append(str(payload))
                elif content_type == "text/html":
                    text_parts.append(_strip_html(str(payload)))
        else:
            payload = message.get_content()
            if isinstance(payload, str):
                text_parts.append(payload)

        text = "\n".join(part for part in text_parts if part).strip()
        return ExtractionResult(
            extracted_text=text or None,
            extraction_quality=_quality_from_text(text),
            extractor_name="eml",
            metadata=metadata,
        )
    except Exception as exc:
        return ExtractionResult(errors=[f"EML extraction failed: {exc}"], extractor_name="eml")


def _extract_msg(file_path: Path) -> ExtractionResult:
    if extract_msg is None:
        return ExtractionResult(errors=["extract-msg not available"], extractor_name="msg")
    try:
        message = extract_msg.Message(str(file_path))
        text_parts = []
        metadata = {
            "subject": message.subject,
            "from": message.sender,
            "to": message.to,
        }
        if message.subject:
            text_parts.append(f"Subject: {message.subject}")
        if message.sender:
            text_parts.append(f"From: {message.sender}")
        if message.to:
            text_parts.append(f"To: {message.to}")
        if message.body:
            text_parts.append(message.body)
        text = "\n".join(part for part in text_parts if part).strip()
        return ExtractionResult(
            extracted_text=text or None,
            extraction_quality=_quality_from_text(text),
            extractor_name="msg",
            metadata=metadata,
        )
    except Exception as exc:
        return ExtractionResult(errors=[f"MSG extraction failed: {exc}"], extractor_name="msg")


EXTRACTORS: Dict[str, Callable[[Path], ExtractionResult]] = {
    ".pdf": _extract_pdf,
    ".docx": _extract_docx,
    ".doc": _extract_doc,
    ".pptx": _extract_pptx,
    ".ppt": _extract_ppt,
    ".xlsx": _extract_xlsx,
    ".xls": _extract_xls,
    ".eml": _extract_eml,
    ".msg": _extract_msg,
    ".txt": _read_text_file,
    ".md": _read_text_file,
    ".csv": _read_text_file,
    ".log": _read_text_file,
    ".json": _read_text_file,
    ".xml": _read_text_file,
    ".html": _read_text_file,
    ".htm": _read_text_file,
    ".jpg": _extract_image,
    ".jpeg": _extract_image,
    ".png": _extract_image,
}


def extract(file_path: Path) -> ExtractionResult:
    extractor = EXTRACTORS.get(file_path.suffix.lower())
    if not extractor:
        return ExtractionResult(
            extractor_name="unsupported",
            errors=[f"Unsupported file type: {file_path.suffix.lower() or '<none>'}"],
        )
    result = extractor(file_path)
    result.metadata.setdefault("mime_type", sniff_mime_type(file_path))
    result.metadata.setdefault("filename", file_path.name)
    return result
