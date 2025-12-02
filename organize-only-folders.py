import os
import shutil
import ollama
import json
import pypdf
import docx
import pytesseract
import time
import functools
import re
import logging
import subprocess
from pptx import Presentation 
from PIL import Image
from pathlib import Path
from typing import Optional, Tuple, Any, List
import shutil as _shutil  # keep name for which checks if needed

# --- 1. AGENT CONFIGURATION ---

# Change these paths to match your setup
# ATTENTION: Use absolute paths to prevent errors

# "Inbox" folder (Where your unorganized files are)
INBOX_DIR = str(Path.home() / "Dropbox" / "_courses")

# "Root" folder (Where categorized subfolders will be created)
ORGANIZED_DIR = str(Path.home() / "OrganizedFiles")

# Your locally installed Ollama model (e.g., 'llama3', 'phi3:medium')
OLLAMA_MODEL = 'gemma3:1b'

# Folder names to ignore when scanning for categories ---
IGNORE_DIRS = ['_Unprocessed', '_Ignored'] 

# --- END OF CONFIGURATION ---

# New helper utilities (retry, sanitizers, unique path)
logging.basicConfig(level=logging.INFO, format="%(asctime)s %(levelname)s: %(message)s")
logger = logging.getLogger(__name__)

# --- NEW V4: Setup Dual Logging ---
def setup_logging():
    """Configures two loggers: one for agent actions and one for API calls."""
    
    # --- Agent Log (agent.log + Console) ---
    # We get the root logger so all modules use it
    agent_logger = logging.getLogger() 
    agent_logger.setLevel(logging.INFO)
    # Remove existing handlers to avoid duplicates if re-run
    agent_logger.handlers = [] 
    
    # Formatter for agent logs
    formatter = logging.Formatter(
        "%(asctime)s %(levelname)-8s: %(message)s", 
        datefmt="%Y-%m-%d %H:%M:%S"
    )

    # File handler for agent.log
    file_handler = logging.FileHandler("agent.log", encoding="utf-8")
    file_handler.setLevel(logging.INFO)
    file_handler.setFormatter(formatter)
    
    # Stream handler for console
    stream_handler = logging.StreamHandler()
    stream_handler.setLevel(logging.INFO)
    stream_handler.setFormatter(formatter)

    agent_logger.addHandler(file_handler)
    agent_logger.addHandler(stream_handler)

    # --- API Log (api.log) ---
    api_logger = logging.getLogger("OllamaAPI")
    api_logger.setLevel(logging.INFO)
    api_logger.propagate = False # Don't send API logs to the root logger

    # File handler for api.log
    api_file_handler = logging.FileHandler("api.log", encoding="utf-8")
    api_file_handler.setLevel(logging.INFO)
    
    # Simple formatter for API log (request/response)
    api_formatter = logging.Formatter(
        "%(asctime)s --- %(message)s", 
        datefmt="%Y-%m-%d %H:%M:%S"
    )
    api_file_handler.setFormatter(api_formatter)
    
    api_logger.addHandler(api_file_handler)

    # Return the loggers to be passed around
    return agent_logger, api_logger
# --- END NEW V4 ---


def retry(exceptions, tries: int = 3, delay: float = 1.0, backoff: int = 2):
    def decorator(f):
        @functools.wraps(f)
        def wrapper(*args, **kwargs):
            # Assumes logger is passed as a keyword argument or globally available
            logger = kwargs.get('logger', logging.getLogger()) 
            _tries, _delay = tries, delay
            while _tries > 1:
                try:
                    return f(*args, **kwargs)
                except exceptions as e:
                    logger.warning("Transient error: %s. Retrying in %.1fs...", e, _delay)
                    time.sleep(_delay)
                    _tries -= 1
                    _delay *= backoff
            return f(*args, **kwargs)
        return wrapper
    return decorator

def sanitize_category(name: str) -> str:
    """Return a safe folder name from the LLM-provided category."""
    if not name or not isinstance(name, str):
        return "_Unprocessed"
    # remove control chars and reserved file chars, collapse whitespace
    safe = re.sub(r'[<>:"/\\|?*\x00-\x1f]+', '', name).strip()
    safe = re.sub(r'\s+', ' ', safe)
    # limit length
    return safe[:100] or "_Unprocessed"

def unique_dest_path(folder: str, filename: str) -> str:
    """Return a non-colliding path inside folder for filename."""
    base, ext = os.path.splitext(filename)
    candidate = filename
    n = 1
    while os.path.exists(os.path.join(folder, candidate)):
        candidate = f"{base}_{n}{ext}"
        n += 1
    return os.path.join(folder, candidate)


# --- UPDATED V4: Added PPT/PPTX Support ---
def extract_text(file_path: str, logger: logging.Logger) -> Optional[str]:
    """
    Extracts raw text from different file types.
    Returns None if the file type is unsupported or an error occurs.
    If a PDF has no extractable text and Tesseract is available, falls back to OCR.
    """
    ext = os.path.splitext(file_path)[1].lower()
    text_parts = []
    try:
        if ext == '.pdf':
            reader = pypdf.PdfReader(file_path)
            for page in reader.pages:
                # page.extract_text() can be None
                page_text = page.extract_text() or ""
                text_parts.append(page_text)
            joined = "\n".join(text_parts).strip()
            if not joined:
                # Try OCR fallback if tesseract exists
                if is_tesseract_available():
                    logger.info("No text in PDF; attempting OCR fallback for %s", file_path)
                    try:
                        # NOTE: This only OCRs the *first page* if it's a multi-page image-only PDF.
                        # A full multi-page OCR would require pdf2image.
                        img_text = pytesseract.image_to_string(Image.open(file_path), lang='eng')
                        return img_text.strip() or None
                    except Exception as e:
                        logger.warning("OCR fallback failed: %s", e)
                        return None
                return None
            return joined

        elif ext == '.docx':
            doc = docx.Document(file_path)
            for para in doc.paragraphs:
                text_parts.append(para.text)
            return "\n".join(text_parts).strip() or None

        elif ext == '.doc':
            logger.info("Found .doc, attempting soffice conversion: %s", file_path)
            out_dir = os.path.dirname(file_path)
            subprocess.check_call(['soffice', '--headless', '--convert-to', 'docx', '--outdir', out_dir, file_path])
            base_name = os.path.splitext(os.path.basename(file_path))[0]
            docx_path = os.path.join(out_dir, base_name + '.docx')
            
            if not os.path.exists(docx_path):
                 logger.warning("soffice failed to convert .doc: %s", file_path)
                 return None

            doc = docx.Document(docx_path)
            for para in doc.paragraphs:
                text_parts.append(para.text)
            joined = "\n".join(text_parts).strip()
            os.remove(docx_path)  # clean up temp file
            return joined or None
        
        # --- NEW V4: Added PPTX Support ---
        elif ext == '.pptx':
            pres = Presentation(file_path)
            for slide in pres.slides:
                for shape in slide.shapes:
                    if not shape.has_text_frame:
                        continue
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            text_parts.append(run.text)
            return "\n".join(text_parts).strip() or None

        # --- NEW V4: Added PPT (legacy) Support ---
        elif ext == '.ppt':
            logger.info("Found .ppt, attempting soffice conversion: %s", file_path)
            out_dir = os.path.dirname(file_path)
            subprocess.check_call(['soffice', '--headless', '--convert-to', 'pptx', '--outdir', out_dir, file_path])
            base_name = os.path.splitext(os.path.basename(file_path))[0]
            pptx_path = os.path.join(out_dir, base_name + '.pptx')

            if not os.path.exists(pptx_path):
                 logger.warning("soffice failed to convert .ppt: %s", file_path)
                 return None

            # Process the converted .pptx
            pres = Presentation(pptx_path)
            for slide in pres.slides:
                for shape in slide.shapes:
                    if not shape.has_text_frame:
                        continue
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            text_parts.append(run.text)
            
            joined = "\n".join(text_parts).strip()
            os.remove(pptx_path)  # clean up temp file
            return joined or None
        
        elif ext in ['.jpg', '.jpeg', '.png']:
            try:
                return pytesseract.image_to_string(Image.open(file_path), lang='eng').strip() or None
            except Exception as e:
                logger.warning("Image OCR failed for %s: %s", file_path, e)
                return None

        elif ext == '.txt':
            try:
                with open(file_path, 'r', encoding='utf-8') as f:
                    return f.read().strip() or None
            except UnicodeDecodeError:
                with open(file_path, 'r', encoding='latin-1') as f:
                    return f.read().strip() or None
        else:
            # Unsupported file type for text extraction
            logger.warning("Unsupported file type for extraction: %s", file_path)
            return None

    except Exception as e:
        logger.warning("Error extracting text from %s: %s", os.path.basename(file_path), e)
        return None


def get_existing_categories(dir_path: str) -> List[str]:
    """
    Scans the target directory and returns a list of existing folder names.
    """
    categories = []
    try:
        for item in os.listdir(dir_path):
            item_path = os.path.join(dir_path, item)
            # Add if it's a directory and not in our ignore list
            if os.path.isdir(item_path) and item not in IGNORE_DIRS:
                categories.append(item)
    except FileNotFoundError:
        # This is normal on the first run
        pass
    except Exception as e:
        # Note: logger might not be configured yet if called early
        print(f"  [!] Error scanning for categories: {e}") 
    return categories


@retry(Exception, tries=3, delay=1, backoff=2)
def call_ollama_chat(model: str, messages: list, fmt: str = 'json', **kwargs) -> Any:
    """Wrap ollama.chat with retries; bubble exceptions on final failure."""
    # Pass logger to retry wrapper
    return ollama.chat(model=model, messages=messages, format=fmt)


# --- UPDATED V4: Added API Logger ---
def get_classification_from_ollama(
    file_content: str,
    filename: str,
    existing_categories: List[str],
    logger: logging.Logger,
    api_logger: logging.Logger
) -> Optional[dict]:
    """
    Sends the text and *existing categories* to Ollama for classification.
    Validates and sanitizes the response.
    """
    
    # Advanced Slicing Strategy
    if len(file_content) > 20000:
        # Take the first 14,500 characters and the last 5,000 characters
        content_preview = file_content[:14500] + "\n... [CONTENT SKIPPED] ...\n" + file_content[-5000:]
    else:
        content_preview = file_content

    # Normalize all whitespace ---
    content_preview = ' '.join(content_preview.split())

    prompt = f"""
    You are an expert file organization agent. 
    Your goal is to be comprehensive, informed, smart and *consistent*. 
    A list of existing categories (folders) is: {str(existing_categories)}

    Filename: "{filename}"

    ##ContentSample##
    
    "{content_preview}"

    ##EndContentSample##

    Your task:
    1.  Analyze the file and its content. First, check if it clearly belongs to one of the **existing categories**. If the list is empty, you must propose a new category.
    2.  If it fits an existing category, use that **exact** category name.
    3.  If it does **not** fit *any* existing category, propose a **new**, single, descriptive category name in CamelCase (e.g., 'ProductManagement', 'Research', 'Presentations', 'UX', 'Platforms', 'Summaries', 'Programming', etc).
    
    Respond ONLY with a valid JSON object like this example:
    {{
      "category": "CategoryName"
    }}
    """

    # --- NEW V4: Log API Request ---
    api_logger.info("REQUEST:\n%s", prompt)

    try:
        logger.info("Contacting Ollama (context folders: %d)...", len(existing_categories))
        response = call_ollama_chat(
            model=OLLAMA_MODEL,
            messages=[{'role': 'user', 'content': prompt}],
            fmt='json',
            logger=logger  # Pass logger for retry logic
        )

        # Extract content from various response shapes (dicts, objects, ChatResponse)
        content = None
        response_content_str = "Invalid response object"

        try:
            if isinstance(response, dict):
                # try both patterns
                if 'message' in response and isinstance(response['message'], dict):
                    content = response['message'].get('content')
                else:
                    content = response.get('content') or response
            else:
                # object-like response (e.g., ollama._types.ChatResponse)
                # try common attributes in order of likelihood
                if hasattr(response, 'message'):
                    msg = getattr(response, 'message')
                    # msg may be dict-like or object with .content
                    if isinstance(msg, dict):
                        content = msg.get('content')
                    elif hasattr(msg, 'content'):
                        content = getattr(msg, 'content')
                if content is None and hasattr(response, 'content'):
                    content = getattr(response, 'content')
            
            # --- NEW V4: Log API Response ---
            response_content_str = str(content) # Store string representation for logging
            api_logger.info("RESPONSE:\n%s", response_content_str)
            
            # If content is a JSON string, try parsing it
            if isinstance(content, str):
                try:
                    parsed = json.loads(content)
                    content = parsed
                except Exception:
                    logger.warning("Ollama content was a non-JSON string.")
                    pass # Leave as string if not JSON, will fail next check

        except Exception as e:
            logger.warning("Failed to normalize Ollama response (%s): %s", type(response), e)
            api_logger.warning("RESPONSE (RAW): %s", str(response))
            content = None

        if not isinstance(content, dict):
            logger.warning("Ollama content is not a JSON object/dict. See api.log for details.")
            return None

        # Validate fields
        raw_category = content.get('category')

        if not raw_category or not isinstance(raw_category, str):
            logger.warning("Invalid category returned by Ollama. See api.log for details.")
            return None

        # Sanitize outputs
        category = sanitize_category(raw_category)

        return {"category": category, "filename": filename}

    except json.JSONDecodeError:
        logger.error("Ollama returned invalid JSON. See api.log for details.")
        return None
    except Exception as e:
        logger.error("Ollama API error: %s", e)
        return None


# --- UPDATED V4: Pass logger ---
def move_file(original_path: str, category: str, filename: str, logger: logging.Logger) -> Optional[str]:
    """
    Move the original file to the new category folder.
    Returns destination path on success, None on failure.
    """
    if not category or not filename:
        category = "_Unprocessed"

    # Sanitize category and filename
    safe_category = sanitize_category(category)
    dest_folder = os.path.join(ORGANIZED_DIR, safe_category)
    os.makedirs(dest_folder, exist_ok=True)
    dest_path = unique_dest_path(dest_folder, filename)

    try:
        shutil.move(original_path, dest_path)
        # Logger is now passed
        logger.info("Moved: %s -> %s/%s", os.path.basename(original_path), safe_category, os.path.basename(dest_path))
        return dest_path
    except Exception as e:
        logger.error("Error moving %s: %s", os.path.basename(original_path), e)
        return None

def is_tesseract_available() -> bool:
    """Return True if tesseract is installed / available to pytesseract."""
    try:
        # prefer pytesseract API if present
        try:
            _ = pytesseract.get_tesseract_version()
            return True
        except Exception:
            # fallback to checking PATH
            return shutil.which("tesseract") is not None
    except Exception:
        return False

# --- UPDATED V4: Pass logger ---
def check_prereqs(logger: logging.Logger) -> bool:
    """Ensure Ollama is reachable and Tesseract presence (recommended)."""
    try:
        ollama.list()
    except Exception as e:
        logger.critical("Could not connect to Ollama: %s", e)
        return False

    if not is_tesseract_available():
        logger.warning("Tesseract not found. OCR fallbacks will be disabled. Install: brew install tesseract")
    
    # Check for soffice (LibreOffice for .doc/.ppt)
    if not shutil.which("soffice"):
         logger.warning("soffice (LibreOffice) not found. .doc and .ppt conversion will fail. Install LibreOffice.")
         
    return True


# --- UPDATED V4: Pass loggers ---
def main(logger: logging.Logger, api_logger: logging.Logger):
    """
    Main function to orchestrate the process.
    """
    logger.info("--- Starting Organization Agent (V4) ---")
    logger.info(f"Inbox: {INBOX_DIR}")
    logger.info(f"Output: {ORGANIZED_DIR}")
    logger.info(f"AI Model: {OLLAMA_MODEL}\n")

    # Ensure main directories exist
    os.makedirs(INBOX_DIR, exist_ok=True)
    os.makedirs(ORGANIZED_DIR, exist_ok=True)

    # Scan for existing categories BEFORE processing files
    logger.info("Scanning for existing categories...")
    existing_categories = get_existing_categories(ORGANIZED_DIR)
    if existing_categories:
        logger.info(f"  [+] Found: {existing_categories}")
    else:
        logger.info("  [*] No existing categories found. Will create new ones.")

    files_to_process = [f for f in os.listdir(INBOX_DIR) if os.path.isfile(os.path.join(INBOX_DIR, f)) and not f.startswith('.')]

    if not files_to_process:
        logger.info("No files found to organize.")
        return

    for filename in files_to_process:
        file_path = os.path.join(INBOX_DIR, filename)
        logger.info(f"\n--- Processing: {filename} ---")
        
        # 1. Extract Text
        content = extract_text(file_path, logger)
        
        if not content or content.isspace():
            logger.warning("Empty content or unsupported file type. Leaving it unprocessed.")
            continue
        
        # 2. Classify with AI
        # Pass the category list to the classification function
        classification_data = get_classification_from_ollama(
            content, 
            filename, 
            existing_categories,
            logger,
            api_logger
        )
        
        # 3. Move File
        if classification_data and isinstance(classification_data, dict):
            category = classification_data.get('category', 'Misc')
            filename = classification_data.get('filename', filename)

            dest_path = move_file(file_path, category, filename, logger)
            
            if dest_path:
                # --- NEW V4: Summary log entry ---
                logger.info(f"SUCCESS: '{filename}' classified as '{category}' and moved to '{dest_path}'")
                if category not in existing_categories and category not in IGNORE_DIRS:
                    logger.info(f"  [+] New category '{category}' added to context.")
                    existing_categories.append(category)
            else:
                 logger.error(f"FAILURE: '{filename}' classified as '{category}' but MOVE FAILED.")

        else:
            logger.error(f"FAILURE: AI classification failed for '{filename}'. Leaving it unprocessed.")

    logger.info("\n--- Organization Complete ---")


# --- UPDATED V4: Call setup_logging and pass loggers ---
if __name__ == "__main__":
    # Setup logging first
    logger, api_logger = setup_logging()
    
    if not check_prereqs(logger):
        logger.critical("\n[CRITICAL] Pre-run checks failed. Fix issues then retry.")
    else:
        main(logger, api_logger)
